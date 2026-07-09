"""Generate output/jul09_status_meeting_vi.pptx from Jul-9 deck content.

Run from the project root:
    uv run python scripts/generate_jul09_pptx.py
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

ROOT = Path(__file__).resolve().parent.parent
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from scripts.generate_jun11_pptx import (
    CREAM,
    FIG_DIR,
    GOLD,
    GOLD_SOFT,
    IVORY,
    MARGIN,
    MIST,
    MUTED,
    SLIDE_H,
    SLIDE_W,
    TEAL,
    TEAL_DARK,
    _card,
    _picture_fit,
    _solid_bg,
    _stat_card,
    _takeaway,
    _text,
)

OUT_PATH = Path("output/jul09_status_meeting_vi.pptx")
TOTAL_SLIDES = 9


def _content_slide9(prs: Presentation, n: int, kicker: str, title: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, CREAM)
    _text(slide, MARGIN, Inches(0.32), Inches(6.0), Inches(0.3),
          kicker.upper(), size=12, bold=True, color=GOLD)
    _text(slide, MARGIN, Inches(0.58), SLIDE_W - 2 * MARGIN, Inches(0.62),
          title, size=27, bold=True, color=TEAL)
    _text(slide, SLIDE_W - Inches(1.7), SLIDE_H - Inches(0.38),
          Inches(1.2), Inches(0.3), f"{n} / {TOTAL_SLIDES}", size=10, color=MUTED,
          align=PP_ALIGN.RIGHT)
    _text(slide, MARGIN, SLIDE_H - Inches(0.38), Inches(5.0), Inches(0.3),
          "Project 8 - Entanglement in Online Choir", size=10, color=MUTED)
    return slide


def s1_title(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, TEAL_DARK)
    _text(slide, MARGIN, Inches(1.65), SLIDE_W - 2 * MARGIN, Inches(0.9),
          "Status Meeting VI", size=48, bold=True, color=IVORY, align=PP_ALIGN.CENTER)
    _text(slide, MARGIN, Inches(2.75), SLIDE_W - 2 * MARGIN, Inches(0.7),
          "Project 8: Entanglement in Online Choir", size=26, color=GOLD_SOFT,
          align=PP_ALIGN.CENTER)
    _text(slide, MARGIN, Inches(3.85), SLIDE_W - 2 * MARGIN, Inches(0.45),
          "SNA-OSN-M Summer 2026  -  Uni Bamberg x Uni Koeln x HSLU",
          size=15, color=MIST, align=PP_ALIGN.CENTER)
    _text(slide, MARGIN, Inches(4.85), SLIDE_W - 2 * MARGIN, Inches(0.4),
          "Presented by Zuraiz, on behalf of the team", size=16, bold=True,
          color=IVORY, align=PP_ALIGN.CENTER)
    _text(slide, MARGIN, Inches(5.45), SLIDE_W - 2 * MARGIN, Inches(0.4),
          "Supervisors: Prof. Janine Hacker - Prof. Peter Gloor",
          size=13, color=MIST, align=PP_ALIGN.CENTER)
    _text(slide, MARGIN, Inches(6.35), SLIDE_W - 2 * MARGIN, Inches(0.4),
          "2026-07-09 - 14:00 CET", size=14, color=GOLD_SOFT, align=PP_ALIGN.CENTER)


def s2_goals_plan(prs: Presentation) -> None:
    slide = _content_slide9(prs, 2, "Recap", "Goals and plan")
    _text(slide, MARGIN, Inches(1.45), SLIDE_W - 2 * MARGIN, Inches(0.6),
          "We are measuring how well online choirs coordinate when singers are not in the same room.",
          size=19, bold=True, color=TEAL)
    rows = [
        ("E(t)", "One coordination score over time from audio, network, and visual signals."),
        ("H1", "Does higher latency reduce choir coordination?"),
        ("H2", "Do influence networks show leadership structure?"),
        ("H3 and plan", "Test visual/body signals when data exists; prepare the report and final dashboard demo."),
    ]
    y = Inches(2.25)
    for title, body in rows:
        _card(slide, MARGIN, y, SLIDE_W - 2 * MARGIN, Inches(0.82), title, [body])
        y += Inches(0.93)
    _takeaway(slide, "Today we connect the project goal to the latest results and the final-presentation plan.")


def s3_iteration_progress(prs: Presentation) -> None:
    slide = _content_slide9(prs, 3, "Progress", "Progress during the last iteration")
    _text(slide, MARGIN, Inches(1.45), SLIDE_W - 2 * MARGIN, Inches(0.55),
          "Since Status Meeting V, we completed the Jul-9 report checkpoint.",
          size=17, bold=True, color=TEAL)
    rows = [
        ("Report draft v1", "Now writes up the H1 result."),
        ("H1", "Ready to use in the report as the main timing result."),
        ("H2", "Cleaner interpretation: weak leadership signal in human choir networks."),
        ("H3", "Open because the needed paired audio-video data is unavailable."),
        ("Status VI", "Slides, script, and Q&A notes are prepared."),
    ]
    y = Inches(2.05)
    for title, body in rows:
        _card(slide, MARGIN, y, SLIDE_W - 2 * MARGIN, Inches(0.72), title, [body])
        y += Inches(0.82)
    _takeaway(slide, "The work moved from finding results to preparing the final presentation and report.")


def s4_report(prs: Presentation) -> None:
    slide = _content_slide9(prs, 4, "Progress", "Report draft v1 is ready for review")
    _card(slide, MARGIN, Inches(1.45), Inches(5.9), Inches(3.6), "What is in the draft",
          [
              "Abstract, hypotheses, data, methods, results, limitations, conclusion.",
              "H1 presented as the headline result.",
              "H2 presented as a measured leadership signal in human datasets.",
              "H3 presented as an open data-availability limitation.",
          ], size=13)
    _stat_card(slide, Inches(7.0), Inches(1.6), Inches(2.0), Inches(1.25), "H1",
               "Supported in onset timing", number_size=28)
    _stat_card(slide, Inches(9.25), Inches(1.6), Inches(2.0), Inches(1.25), "H2",
               "Partially supported", number_size=28)
    _stat_card(slide, Inches(11.5), Inches(1.6), Inches(1.3), Inches(1.25), "H3",
               "Open", number_size=25)
    _text(slide, Inches(7.0), Inches(3.25), Inches(5.8), Inches(1.6),
          "The Jul-9 hard milestone from Status Meeting V is met: the latency result is in prose, "
          "with methods and limitations already stated.",
          size=15, color=TEAL)
    _takeaway(slide, "The Jul-9 report milestone from Status Meeting V is met.")


def s5_h1(prs: Presentation) -> None:
    slide = _content_slide9(prs, 5, "Result", "H1: latency breaks attack timing")
    fig = FIG_DIR / "tier3_corpus_summary.png"
    if fig.exists():
        _picture_fit(slide, fig, MARGIN, Inches(1.35), Inches(8.2), Inches(4.75))
    _text(slide, Inches(9.15), Inches(1.55), Inches(3.65), Inches(4.3),
          "Result summary:\n\n"
          "- Onset synchrony falls strongly as jitter rises.\n"
          "- Loudness-envelope coupling stays almost flat.\n"
          "- Latency breaks when singers land notes, not how loud they are.\n\n"
          "Clean -> Zoom drop:\n"
          "Dagstuhl 57%\nESMUC 66%\nChoralSynth 76%",
          size=13, color=TEAL)
    _takeaway(slide, "An envelope-only metric would have missed the real latency effect.")


def s6_h2(prs: Presentation) -> None:
    slide = _content_slide9(prs, 6, "Result", "H2: leadership appears in human choir networks")
    _card(slide, MARGIN, Inches(1.45), Inches(5.9), Inches(2.25), "What the original H2 asked",
          [
              "Original: networks become leader-dominated as latency rises.",
              "Current data cannot test that cleanly because injected delay on pre-recorded audio cannot create a behavioral leader.",
          ], size=13)
    _card(slide, MARGIN, Inches(3.95), Inches(5.9), Inches(1.7), "Current H2 result",
          [
              "Clean human choir influence networks carry weak leadership structure above a density-matched random null."
          ], size=13)
    _stat_card(slide, Inches(7.0), Inches(1.65), Inches(1.7), Inches(1.25), "3/5",
               "Dagstuhl significant", number_size=28)
    _stat_card(slide, Inches(9.0), Inches(1.65), Inches(1.7), Inches(1.25), "2/3",
               "ESMUC significant", number_size=28)
    _stat_card(slide, Inches(11.0), Inches(1.65), Inches(1.7), Inches(1.25), "2/20",
               "ChoralSynth, approx. chance", number_size=25)
    _text(slide, Inches(7.0), Inches(3.35), Inches(5.7), Inches(1.5),
          "Interpretation: leadership appears as a weak human coordination signal, "
          "not as a synthetic-rendering artifact. We do not claim strong hierarchy.",
          size=14, color=TEAL)
    _text(slide, MARGIN, Inches(5.55), SLIDE_W - 2 * MARGIN, Inches(0.85),
          "Definition (operational): leader dominance = out-degree centralization of the directed "
          "influence graph, measured as the Gini coefficient of node out-degree. 0 = every singer "
          "exerts equal outgoing influence (democratic); toward 1 = one singer Granger-causes the "
          "others without following back (a leader). Test: observed Gini vs 1000 random graphs of "
          "identical size and density. Observed corpus mean 0.154 vs null 0.139.",
          size=11.5, color=MUTED)
    _takeaway(slide, "Leadership appears as a human coordination signal, not a synthetic-rendering artifact.")


def s7_dashboard(prs: Presentation) -> None:
    slide = _content_slide9(prs, 7, "Progress", "Dashboard alpha uses real outputs")
    fig = FIG_DIR / "wp4_dashboard_realdata.png"
    if fig.exists():
        _picture_fit(slide, fig, MARGIN, Inches(1.35), Inches(8.2), Inches(4.75))
    _text(slide, Inches(9.1), Inches(1.55), Inches(3.7), Inches(4.25),
          "Real alpha:\n\n"
          "- E(t) from the real pipeline\n"
          "- Graph from real Granger GEXF\n"
          "- Pose overlay from real parquet\n"
          "- Metadata shows available signals\n\n"
          "This screenshot shows the current real-data alpha used for demo planning.\n\n"
          "Demo structure:\n"
          "One audio/network example shows E(t) and the influence graph. One video/pose example shows playback and pose overlay.",
          size=12, color=TEAL)
    _takeaway(slide, "The final demo can be clear and honest even though no current piece has all three signals.")


def s8_remaining(prs: Presentation) -> None:
    slide = _content_slide9(prs, 8, "Next", "Plan for the next iteration")
    _text(slide, MARGIN, Inches(1.28), SLIDE_W - 2 * MARGIN, Inches(0.35),
          "The next iteration is the final presentation and report phase.",
          size=16, bold=True, color=TEAL)
    rows = [
        ("Final presentation", "Build the 20-minute Jul-23 deck and demo narration."),
        ("Report", "Convert draft v1 into the final 10-20 page seminar report."),
        ("Reproducibility", "Run final checks and regenerate headline figures from committed outputs."),
        ("Dashboard", "Prepare the presentation laptop and rehearse the demo path."),
        ("Limitations", "Present H3 and latency-driven H2 as open questions for future data."),
    ]
    y = Inches(1.78)
    for title, body in rows:
        _card(slide, MARGIN, y, SLIDE_W - 2 * MARGIN, Inches(0.82), title, [body])
        y += Inches(0.93)
    _takeaway(slide, "The remaining work is final packaging, verification, and rehearsal.")


def s9_final_alignment(prs: Presentation) -> None:
    slide = _content_slide9(prs, 9, "Alignment", "What we need to align before the final presentation")
    col = Inches(3.95)
    gap = Inches(0.18)
    _card(slide, MARGIN, Inches(1.45), col, Inches(3.8), "Retrospective output",
          [
              "Null results were kept and explained.",
              "Constant-delay failure led to the stronger jitter/onset result.",
              "Dataset claims now trace to files and generated outputs.",
          ], size=13)
    _card(slide, MARGIN + col + gap, Inches(1.45), col, Inches(3.8), "Proposed structure",
          [
              "Research question and hypotheses.",
              "Method, datasets, results, and interpretation.",
              "Implementation, limitations, and next steps.",
          ], size=13)
    _card(slide, MARGIN + 2 * (col + gap), Inches(1.45), col, Inches(3.8), "Questions for feedback",
          [
              "Final structure: research question, hypotheses, method, results, implementation, limitations, next steps.",
              "Main focus: what should receive the most attention in the final presentation?",
              "Dashboard format: is a screenshot enough, or should we show it live?",
          ], size=13)
    _text(slide, MARGIN, Inches(5.75), SLIDE_W - 2 * MARGIN, Inches(0.45),
          "Thank you.", size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


def build() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    for fn in (
        s1_title,
        s2_goals_plan,
        s3_iteration_progress,
        s4_report,
        s5_h1,
        s6_h2,
        s7_dashboard,
        s8_remaining,
        s9_final_alignment,
    ):
        fn(prs)
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"Wrote {OUT_PATH} ({len(prs.slides)} slides, {OUT_PATH.stat().st_size / 1024:.0f} KB)")


if __name__ == "__main__":
    build()
