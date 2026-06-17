"""Generate output/jun25_status_meeting_v.pptx from jun25_deck.md content.

Reuses the Jun-11 generator's design helpers (palette, kicker/title/footer,
stat cards, takeaway banner, aspect-correct figure fit) for visual
consistency; only the 8 slides' content differs.

Run from the project root:
    uv run python scripts/generate_jun25_pptx.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from scripts.generate_jun11_pptx import (
    FIG_DIR,
    GOLD_SOFT,
    IVORY,
    MARGIN,
    MIST,
    SLIDE_H,
    SLIDE_W,
    TEAL,
    TEAL_DARK,
    _card,
    _content_slide,
    _picture_fit,
    _solid_bg,
    _stat_card,
    _takeaway,
    _text,
)

OUT_PATH = Path("output/jun25_status_meeting_v.pptx")


def s1_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, TEAL_DARK)
    _text(slide, MARGIN, Inches(1.7), SLIDE_W - 2 * MARGIN, Inches(0.9),
          "Status Meeting V", size=48, bold=True, color=IVORY, align=PP_ALIGN.CENTER)
    _text(slide, MARGIN, Inches(2.8), SLIDE_W - 2 * MARGIN, Inches(0.7),
          "Project 8: Entanglement in Online Choir", size=26, color=GOLD_SOFT, align=PP_ALIGN.CENTER)
    _text(slide, MARGIN, Inches(4.0), SLIDE_W - 2 * MARGIN, Inches(0.45),
          "SNA-OSN-M Summer 2026  ·  Uni Bamberg × Uni Köln × HSLU",
          size=15, color=MIST, align=PP_ALIGN.CENTER)
    _text(slide, MARGIN, Inches(5.0), SLIDE_W - 2 * MARGIN, Inches(0.4),
          "Supervisors: Prof. Janine Hacker · Prof. Peter Gloor", size=13, color=MIST,
          align=PP_ALIGN.CENTER)
    _text(slide, MARGIN, Inches(6.4), SLIDE_W - 2 * MARGIN, Inches(0.4),
          "2026-06-25 · 14:00 CET", size=14, color=GOLD_SOFT, align=PP_ALIGN.CENTER)


def s2_recap(prs):
    slide = _content_slide(prs, 2, "Recap", "Goals and plan")
    _text(slide, MARGIN, Inches(1.5), SLIDE_W - 2 * MARGIN, Inches(0.5),
          "E(t) = mean( A audio, V visual, N network ): one coordination score per moment",
          size=16, bold=True, color=TEAL)
    h = Inches(1.0)
    _card(slide, MARGIN, Inches(2.2), SLIDE_W - 2 * MARGIN, h, "H1 · Latency",
          ["Lower latency → higher coordination. (This iteration's focus.)"])
    _card(slide, MARGIN, Inches(3.3), SLIDE_W - 2 * MARGIN, h, "H2 · Topology",
          ["Influence network shifts democratic → leader-dominated as latency rises."])
    _card(slide, MARGIN, Inches(4.4), SLIDE_W - 2 * MARGIN, h, "H3 · Visual signals",
          ["Body sway + breathing add information over audio (still data-blocked)."])
    _text(slide, MARGIN, Inches(5.7), SLIDE_W - 2 * MARGIN, Inches(0.7),
          "Status meeting 5 of 6. Jun 11: E(t) operational. Today: first test of H1 via latency injection.",
          size=13, color=TEAL)
    _takeaway(slide, "This iteration: we found the latency signal, and learned exactly where it lives.")


def s3_headline(prs):
    slide = _content_slide(prs, 3, "Progress · last iteration",
                           "We found the H1 latency signal (in attack timing)")
    steps = [
        ("1 · Constant delay", "No E(t) effect. Our coupling is lag-tolerant; it absorbs a constant shift."),
        ("2 · Realistic jitter", "Still flat. 10s loudness envelopes are robust to tens-of-ms timing noise."),
        ("3 · Onset synchrony", "Zero-lag attack timing, the quantity latency breaks. The signal appears."),
    ]
    y = Inches(1.5)
    for t, b in steps:
        _card(slide, MARGIN, y, SLIDE_W - 2 * MARGIN, Inches(1.05), t, [b])
        y += Inches(1.15)
    _text(slide, MARGIN, Inches(5.1), SLIDE_W - 2 * MARGIN, Inches(1.3),
          "Result: attack-timing synchrony falls 57-76% (clean → Zoom) across 28 pieces; "
          "loudness coupling barely moves. A loudness-only metric would have said "
          "'latency doesn't matter', and been wrong.",
          size=14, color=TEAL)
    _takeaway(slide, "Latency degrades attack timing, not loudness coupling. The dissociation is the finding.")


def s4_figure(prs):
    slide = _content_slide(prs, 4, "Progress · last iteration", "The dissociation, in one figure")
    fig = FIG_DIR / "tier3_corpus_summary.png"
    if fig.exists():
        _picture_fit(slide, fig, MARGIN, Inches(1.4), SLIDE_W - 2 * MARGIN, Inches(4.6))
    _text(slide, MARGIN, Inches(6.0), SLIDE_W - 2 * MARGIN, Inches(0.7),
          "Left: onset synchrony falls with jitter (all 3 datasets). Right: envelope E(t) flat. "
          "Jitter SDs are measured, not tuned; the constant-delay dead end is reported, not hidden.",
          size=12, color=TEAL)
    _takeaway(slide, "Same recordings, two measures, opposite verdicts.")


def s5_replication(prs):
    slide = _content_slide(prs, 5, "Progress · last iteration", "It replicates across 3 datasets")
    rail = Inches(3.4)
    _stat_card(slide, MARGIN, Inches(1.6), rail, Inches(1.5), "−57%", "Dagstuhl (real human, 5)")
    _stat_card(slide, MARGIN + rail + Inches(0.2), Inches(1.6), rail, Inches(1.5), "−66%", "ESMUC (real human, 3)")
    _stat_card(slide, MARGIN + 2 * (rail + Inches(0.2)), Inches(1.6), rail, Inches(1.5), "−76%", "ChoralSynth (synthetic, 20)")
    _text(slide, MARGIN, Inches(3.5), SLIDE_W - 2 * MARGIN, Inches(1.2),
          "Onset-synchrony drop clean → Zoom. 28 pieces, two independent human datasets plus "
          "synthetic, all monotonic. Both Tier-2 datasets added this iteration via one unified "
          "data adapter (downloaded + MD5-verified vs Zenodo).",
          size=14, color=TEAL)
    _takeaway(slide, "The effect holds across real and synthetic choirs, not a single-dataset artifact.")


def s6_dashboard(prs):
    slide = _content_slide(prs, 6, "Progress · last iteration", "Dashboard alpha, on real data")
    fig = FIG_DIR / "wp4_dashboard_scaffold.png"
    if fig.exists():
        _picture_fit(slide, fig, MARGIN, Inches(1.4), Inches(8.2), Inches(4.7))
    _text(slide, Inches(9.0), Inches(1.6), Inches(3.8), Inches(4.4),
          "Now on real outputs:\n\n• Timeline = real E(t)\n• Graph = real Granger\n  who-leads-whom\n• Video + 33-point pose\n  overlay synced to play\n• Each piece shows the\n  signals it has\n\nMeets the Jun-21 alpha\nmilestone.",
          size=13, color=TEAL)
    _takeaway(slide, "The dashboard runs on the pipeline's real outputs, not mock data.")


def s7_next(prs):
    slide = _content_slide(prs, 7, "Next iteration", "Plan: Jun 26 → Jul 9")
    rows = [
        ("Integration", "Fold onset synchrony into E(t) as its timing-sensitive component"),
        ("WP3", "Finish cross-dataset corpus; per-window time-varying networks"),
        ("WP2", "Pose on all 29 Tier-1 videos (visibility-triaged)"),
        ("Paper", "Draft report methods + results (the latency finding)"),
        ("Compute", "Cluster access for the paper-scale run (2000-shuffle null)"),
    ]
    y = Inches(1.6)
    for track, work in rows:
        _card(slide, MARGIN, y, SLIDE_W - 2 * MARGIN, Inches(0.82), track, [work])
        y += Inches(0.92)
    _takeaway(slide, "From 'found the signal' to 'written up and corpus-complete'.")


def s8_retro(prs):
    slide = _content_slide(prs, 8, "Retrospective · questions", "Retrospective and open questions")
    col = Inches(3.95)
    _card(slide, MARGIN, Inches(1.5), col, Inches(3.4), "What worked",
          ["A wrong first method (constant delay) was caught by its own control and became a sharper one.",
           "", "Every dataset claim traces to a verified file."], size=12.5)
    _card(slide, MARGIN + col + Inches(0.18), Inches(1.5), col, Inches(3.4), "Limitations",
          ["Injection tests transmission timing, not a live singer's adaptation.",
           "", "Envelope E(t) alone is latency-blind (why onset synchrony matters)."], size=12.5)
    _card(slide, MARGIN + 2 * (col + Inches(0.18)), Inches(1.5), col, Inches(3.4), "Questions",
          ["Hacker: foreground onset synchrony over composite E(t) in the report?",
           "", "Coordinators: cluster access for the paper-scale run (ki-support contacted)?"], size=12.5)
    _text(slide, MARGIN, Inches(5.3), SLIDE_W - 2 * MARGIN, Inches(0.5),
          "Thank you.", size=20, bold=True, color=TEAL, align=PP_ALIGN.CENTER)


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    for fn in (s1_title, s2_recap, s3_headline, s4_figure, s5_replication,
               s6_dashboard, s7_next, s8_retro):
        fn(prs)
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"Wrote {OUT_PATH} ({len(prs.slides)} slides, {OUT_PATH.stat().st_size/1024:.0f} KB)")


if __name__ == "__main__":
    build()
