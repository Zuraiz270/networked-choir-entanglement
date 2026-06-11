"""Generate output/jun11_status_meeting_iv.pptx from the jun11_deck.md content.

8 slides following the coordinators' rubric: goals+plan recap, progress
(3 slides), next-iteration plan, retrospective, problems/questions.

Design language (v2, after the v1 text-wall was rejected):
- figure-led slides: the figure is the hero, text is the rail
- stat cards (rounded rectangles, big number + small label) over bullets
- kicker tag above each title locating the slide in the rubric
- footer with slide number on every content slide
- Studio Acoustic palette (teal + gold) for project-family consistency

Run from the project root:
    uv run python scripts/generate_jun11_pptx.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUT_PATH = ROOT / "output" / "jun11_status_meeting_iv.pptx"
FIG_DIR = ROOT / "data" / "figures"

# --- Studio Acoustic palette ---
TEAL = RGBColor(0x0E, 0x4D, 0x5E)
TEAL_DARK = RGBColor(0x07, 0x33, 0x40)
TEAL_MID = RGBColor(0x32, 0x80, 0x95)
GOLD = RGBColor(0xC4, 0x90, 0x2A)
GOLD_SOFT = RGBColor(0xE6, 0xC6, 0x70)
CREAM = RGBColor(0xFA, 0xF6, 0xEC)
IVORY = RGBColor(0xFD, 0xFB, 0xF5)
CHARCOAL = RGBColor(0x24, 0x26, 0x2E)
MUTED = RGBColor(0x6E, 0x6A, 0x60)
MIST = RGBColor(0xCB, 0xD9, 0xDD)
GREEN = RGBColor(0x4F, 0x8E, 0x5A)
RED_SOFT = RGBColor(0xB5, 0x5A, 0x44)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.55)

# Segoe UI ships with Windows; reads cleaner than the Calibri theme default.
FONT_HEAD = "Segoe UI Semibold"
FONT_BODY = "Segoe UI"


def _no_midword(p):
    """Forbid Latin mid-word line breaks (PowerPoint otherwise splits
    words at the shape border: 'th e', 'pa ssing', ...)."""
    pPr = p._p.get_or_add_pPr()
    pPr.set("latinLnBrk", "0")


def _style(run, *, size, bold=False, color=CHARCOAL, head=False):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = FONT_HEAD if head else FONT_BODY


def _solid_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _text(slide, left, top, width, height, text, *, size=16, bold=False,
          color=CHARCOAL, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
          line_spacing=1.0, head=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    _no_midword(p)
    run = p.add_run()
    run.text = text
    _style(run, size=size, bold=bold, color=color, head=head)
    return tb


def _multiline(slide, left, top, width, height, items, *, size=14,
               color=CHARCOAL, gap=6):
    """Paragraph list without bullet glyph clutter. items: list of (text, color|None, bold)."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    for i, item in enumerate(items):
        text, item_color, bold = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.alignment = PP_ALIGN.LEFT
        _no_midword(p)
        run = p.add_run()
        run.text = text
        _style(run, size=size, bold=bold, color=item_color or color)
    return tb


def _kicker(slide, text):
    _text(slide, MARGIN, Inches(0.32), Inches(6.0), Inches(0.3),
          text.upper(), size=12, bold=True, color=GOLD, head=True)


def _title(slide, text):
    _text(slide, MARGIN, Inches(0.58), SLIDE_W - 2 * MARGIN, Inches(0.62),
          text, size=27, bold=True, color=TEAL, head=True)


def _takeaway(slide, text):
    """Teal banner pinned at the slide bottom; one factual line."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN, Inches(6.72),
        SLIDE_W - 2 * MARGIN, Inches(0.52))
    bar.adjustments[0] = 0.06
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _no_midword(p)
    run = p.add_run()
    run.text = text
    _style(run, size=14, bold=True, color=IVORY, head=True)


def _footer(slide, n):
    _text(slide, SLIDE_W - Inches(1.7), SLIDE_H - Inches(0.38),
          Inches(1.2), Inches(0.3), f"{n} / 8", size=10, color=MUTED,
          align=PP_ALIGN.RIGHT)
    _text(slide, MARGIN, SLIDE_H - Inches(0.38), Inches(5.0), Inches(0.3),
          "Project 8 · Entanglement in Online Choir", size=10, color=MUTED)


def _stat_card(slide, left, top, width, height, number, label,
               *, number_color=TEAL, number_size=30):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.adjustments[0] = 0.06
    card.fill.solid()
    card.fill.fore_color.rgb = IVORY
    card.line.color.rgb = MIST
    card.line.width = Pt(1.0)
    tf = card.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    _no_midword(p1)
    r1 = p1.add_run()
    r1.text = number
    _style(r1, size=number_size, bold=True, color=number_color, head=True)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    _no_midword(p2)
    r2 = p2.add_run()
    r2.text = label
    _style(r2, size=11.5, color=MUTED)


def _card(slide, left, top, width, height, title, lines, *, title_color=TEAL,
          fill=IVORY, size=12.5):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.adjustments[0] = 0.04
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = MIST
    card.line.width = Pt(1.0)
    tf = card.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.22)
    tf.margin_right = Inches(0.22)
    tf.margin_top = Inches(0.16)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    _no_midword(p)
    r = p.add_run()
    r.text = title
    _style(r, size=14, bold=True, color=title_color, head=True)
    p.space_after = Pt(6)
    for line in lines:
        lp = tf.add_paragraph()
        lp.alignment = PP_ALIGN.LEFT
        lp.space_after = Pt(4)
        _no_midword(lp)
        lr = lp.add_run()
        lr.text = line
        _style(lr, size=size, color=CHARCOAL)


def _picture_fit(slide, path, left, top, max_w, max_h):
    """Insert picture scaled to fit the (max_w, max_h) box, preserving ratio."""
    from PIL import Image  # pillow ships with mediapipe deps

    with Image.open(path) as im:
        w_px, h_px = im.size
    ratio = min(max_w / w_px, max_h / h_px)
    w, h = int(w_px * ratio), int(h_px * ratio)
    # center within the box
    x = left + (max_w - w) // 2
    y = top + (max_h - h) // 2
    slide.shapes.add_picture(str(path), x, y, width=w, height=h)


def _content_slide(prs, n, kicker, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, CREAM)
    _kicker(slide, kicker)
    _title(slide, title)
    _footer(slide, n)
    return slide


# ---------------- Slides ----------------


def slide_1_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, TEAL_DARK)
    _text(slide, MARGIN, Inches(1.55), SLIDE_W - 2 * MARGIN, Inches(0.9),
          "Status Meeting IV", size=48, bold=True, color=IVORY, align=PP_ALIGN.CENTER)
    _text(slide, MARGIN, Inches(2.65), SLIDE_W - 2 * MARGIN, Inches(0.7),
          "Project 8: Entanglement in Online Choir", size=27, color=GOLD_SOFT,
          align=PP_ALIGN.CENTER)
    line = slide.shapes.add_connector(1, Inches(5.6), Inches(3.55), Inches(7.7), Inches(3.55))
    line.line.color.rgb = GOLD
    line.line.width = Pt(2.5)
    _text(slide, MARGIN, Inches(3.85), SLIDE_W - 2 * MARGIN, Inches(0.45),
          "SNA-OSN-M Summer 2026  ·  Uni Bamberg × Uni Köln × HSLU",
          size=15, color=MIST, align=PP_ALIGN.CENTER)
    _text(slide, MARGIN, Inches(4.95), SLIDE_W - 2 * MARGIN, Inches(0.5),
          "Presented by Hassan Ahmed, on behalf of the team", size=18, bold=True,
          color=IVORY, align=PP_ALIGN.CENTER)
    _text(slide, MARGIN, Inches(5.55), SLIDE_W - 2 * MARGIN, Inches(0.4),
          "Supervisors: Prof. Janine Hacker (Uni Bamberg) · Prof. Peter Gloor (MIT/Köln)",
          size=13, color=MIST, align=PP_ALIGN.CENTER)
    _text(slide, MARGIN, Inches(6.45), SLIDE_W - 2 * MARGIN, Inches(0.4),
          "2026-06-11 · 14:00 CET", size=14, color=GOLD_SOFT, align=PP_ALIGN.CENTER)


def slide_2_recap(prs):
    slide = _content_slide(prs, 2, "Recap", "Goals and plan")

    # Left: formula card + hypotheses
    formula = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN, Inches(1.45), Inches(6.1), Inches(1.05))
    formula.adjustments[0] = 0.08
    formula.fill.solid()
    formula.fill.fore_color.rgb = TEAL
    formula.line.fill.background()
    tf = formula.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _no_midword(p)
    r = p.add_run()
    r.text = "E(t)  =  mean(  A(t) audio,  V(t) visual,  N(t) network  )"
    _style(r, size=19, bold=True, color=IVORY, head=True)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    _no_midword(p2)
    r2 = p2.add_run()
    r2.text = "a single coordination score per time window"
    _style(r2, size=12, color=MIST)

    hyp_y = Inches(2.75)
    hyp_h = Inches(0.98)
    _card(slide, MARGIN, hyp_y, Inches(6.1), hyp_h, "H1 · Latency regimes",
          ["Low-latency tools (Jamulus, SoundJack) score higher E(t) than Zoom."])
    _card(slide, MARGIN, hyp_y + hyp_h + Inches(0.12), Inches(6.1), hyp_h,
          "H2 · Network topology",
          ["Influence network shifts democratic → leader-dominated as latency rises."])
    _card(slide, MARGIN, hyp_y + 2 * (hyp_h + Inches(0.12)), Inches(6.1), hyp_h,
          "H3 · Honest signals",
          ["Body sway + breathing add ≥ 10 points of explained variance over audio."])

    # Right: timeline strip (vertical)
    tl_x = Inches(7.2)
    tl_w = Inches(5.5)
    _text(slide, tl_x, Inches(1.45), tl_w, Inches(0.35),
          "Where we are", size=15, bold=True, color=TEAL)
    milestones = [
        ("Apr 15-16", "Block course + status #1", True, False),
        ("Apr 30", "Status #2 · goals + plan", True, False),
        ("May 21", "Status #3 · all 4 WPs running + Virtual Mirror", True, False),
        ("Jun 11", "Status #4 · TODAY: E(t) operational", False, True),
        ("Jun 25", "Status #5 · cross-regime results", False, False),
        ("Jul 9", "Status #6 · last check-in before final", False, False),
        ("Jul 23 / 31", "Final presentation / final paper", False, False),
    ]
    y = Inches(1.9)
    row_h = Inches(0.62)
    for date, label, done, today in milestones:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, tl_x, y + Inches(0.07),
                                     Inches(0.22), Inches(0.22))
        dot.fill.solid()
        dot.fill.fore_color.rgb = GREEN if done else (GOLD if today else MIST)
        dot.line.fill.background()
        _text(slide, tl_x + Inches(0.38), y, Inches(1.25), Inches(0.5),
              date, size=11.5, bold=True, color=TEAL if not today else GOLD)
        _text(slide, tl_x + Inches(1.7), y, tl_w - Inches(1.7), Inches(0.55),
              label, size=11.5, color=CHARCOAL if not today else GOLD,
              line_spacing=0.95)
        y += row_h

    _takeaway(slide, "Goals unchanged since April. This iteration delivered the first end-to-end E(t).")


def slide_3_headline(prs):
    slide = _content_slide(prs, 3, "Progress · last iteration",
                           "E(t) is operational: all 5 pieces significant above the null")
    fig = FIG_DIR / "et_corpus_comparison.png"
    if fig.exists():
        _picture_fit(slide, fig, MARGIN, Inches(1.4), Inches(8.6), Inches(5.1))

    rail_x = Inches(9.5)
    rail_w = Inches(3.3)
    card_h = Inches(1.18)
    _stat_card(slide, rail_x, Inches(1.5), rail_w, card_h, "5 / 5",
               "pieces significantly above the null (p < 0.005)")
    _stat_card(slide, rail_x, Inches(2.85), rail_w, card_h, "200",
               "circular-shift permutations per piece")
    _stat_card(slide, rail_x, Inches(4.2), rail_w, card_h, "0.57–0.80",
               "observed mean E(t) across pieces", number_size=24)
    _text(slide, rail_x, Inches(5.5), rail_w, Inches(1.1),
          "Split follows the music: homophonic chant scores high, polyphony lower. "
          "The metric reacts to what the choir is doing.",
          size=12, color=MUTED, line_spacing=1.05)

    _takeaway(slide, "E(t) is implemented, reproducible, and significant against the null on all five pieces.")


def slide_4_audio_network(prs):
    slide = _content_slide(prs, 4, "Progress · last iteration",
                           "Audio and network pipelines behind E(t)")
    fig = FIG_DIR / "wp3_influence_graphs_5pieces.png"
    if fig.exists():
        _picture_fit(slide, fig, MARGIN, Inches(1.4), Inches(8.2), Inches(5.1))

    rail_x = Inches(9.1)
    rail_w = Inches(3.7)
    card_h = Inches(1.0)
    _stat_card(slide, rail_x, Inches(1.45), rail_w, card_h, "25",
               "Dagstuhl takes processed (was 1 last sprint)")
    _stat_card(slide, rail_x, Inches(2.6), rail_w, card_h, "288",
               "pairwise audio couplings computed")
    _stat_card(slide, rail_x, Inches(3.75), rail_w, card_h, "2",
               "causality methods per piece (Granger + COP-GC)")
    _text(slide, rail_x, Inches(4.95), rail_w, Inches(1.6),
          "Methods agree on quartets, diverge on full choir (42 vs 25 of 56 edges). "
          "That gap is itself a finding: about 40% of standard edges (17 of 42) "
          "rely on linear magnitude, not pattern structure.",
          size=12, color=MUTED, line_spacing=1.05)

    _takeaway(slide, "Sprint-2 reference reproduces (11/12 edges, density 0.917); the pipeline now scales to the corpus.")


def slide_5_video_dashboard(prs):
    slide = _content_slide(prs, 5, "Progress · last iteration",
                           "Video features and dashboard scaffold")
    fig = FIG_DIR / "wp4_dashboard_scaffold.png"
    if fig.exists():
        _picture_fit(slide, fig, MARGIN, Inches(1.4), Inches(6.4), Inches(5.1))

    right_x = Inches(7.15)
    right_w = Inches(5.6)

    mini = FIG_DIR / "wp2_visual_features_v2.png"
    if mini.exists():
        _picture_fit(slide, mini, right_x, Inches(1.4), right_w, Inches(2.25))

    # three stat cards in one row under the mini figure (no overlaps)
    card_w = Inches(1.78)
    card_gap = Inches(0.13)
    card_y = Inches(3.85)
    _stat_card(slide, right_x, card_y, card_w, Inches(1.0), "10",
               "videos pose-processed", number_size=24)
    _stat_card(slide, right_x + card_w + card_gap, card_y, card_w, Inches(1.0), "5/10",
               "pass 50% detection", number_size=24)
    _stat_card(slide, right_x + 2 * (card_w + card_gap), card_y, card_w, Inches(1.0), "23/23",
               "tests green", number_size=24)

    _multiline(slide, right_x, Inches(5.05), right_w, Inches(1.55),
               [
                   ("WP2: 5 of 10 stratified videos have usable pose tracks (best 98.5% detection); the rest are UI captures or low-res tile grids. The five passing videos form the working set.", None, False),
                   ("WP4: the scaffold renders all four panels end-to-end; real-data wiring follows next iteration.", None, False),
               ], size=12, gap=5)

    _takeaway(slide, "All four work packages advanced; no blockers.")


def slide_6_next_iteration(prs):
    slide = _content_slide(prs, 6, "Next iteration", "Plan: Jun 12 → Jun 25")
    rows = [
        ("WP1 audio", "Per-window Granger → time-varying N(t) for the dashboard timeline"),
        ("WP2 video", "Pose on remaining Tier-1 videos, quality-first triage"),
        ("WP3 network", "Tier-3 latency injection: synthetic jitter at 4 regime levels, E(t) per level. First cross-regime test of H1 + H2."),
        ("WP4 dashboard", "Swap mock JSON for real parquet readers + pose overlay"),
        ("Data", "Download ChoralSynth (Zenodo, research licence) · follow up ESMUC"),
    ]
    t = slide.shapes.add_table(len(rows) + 1, 2, MARGIN, Inches(1.5),
                               SLIDE_W - 2 * MARGIN, Inches(4.3)).table
    t.columns[0].width = Inches(2.6)
    t.columns[1].width = Inches(9.6)
    header = ("Track", "What ships")
    for c, htext in enumerate(header):
        cell = t.cell(0, c)
        cell.text = htext
        cell.fill.solid()
        cell.fill.fore_color.rgb = TEAL
        for para in cell.text_frame.paragraphs:
            _no_midword(para)
            for run in para.runs:
                _style(run, size=14, bold=True, color=IVORY, head=True)
    for r, (track, work) in enumerate(rows, start=1):
        for c, txt in enumerate((track, work)):
            cell = t.cell(r, c)
            cell.text = txt
            cell.fill.solid()
            cell.fill.fore_color.rgb = IVORY if r % 2 else CREAM
            for para in cell.text_frame.paragraphs:
                _no_midword(para)
                for run in para.runs:
                    _style(run, size=13.5, bold=c == 0,
                           color=TEAL if c == 0 else CHARCOAL, head=c == 0)

    _text(slide, MARGIN, Inches(6.0), SLIDE_W - 2 * MARGIN, Inches(0.5),
          "Hard milestone before status #5: dashboard alpha on real data + first Tier-3 cross-regime result.",
          size=15, bold=True, color=GOLD)

    _takeaway(slide, "Next iteration: from measuring coordination to comparing latency regimes.")


def slide_7_retro(prs):
    slide = _content_slide(prs, 7, "Retrospective", "Sprint 3 retrospective")
    col_w = Inches(3.95)
    col_h = Inches(4.35)
    gap = Inches(0.18)
    x0 = MARGIN
    nbh = "‑"  # non-breaking hyphen: keeps "Tier-2" etc. on one line

    _card(slide, x0, Inches(1.5), col_w, col_h, "What worked",
          [
              "One reviewable artefact per work package per iteration.",
              "",
              "Documentation updated at every milestone; project state is readable from three files.",
          ], title_color=GREEN, size=13.5)
    _card(slide, x0 + col_w + gap, Inches(1.5), col_w, col_h, "What went wrong",
          [
              f"ESMUC and ChoralSynth not yet in Tier{nbh}2. ChoralSynth is freely available on Zenodo for research and scheduled for next iteration; ESMUC requires a license (open question).",
              "",
              f"Half of the Tier{nbh}1 videos are screen captures without visible singers. Future curation will filter on singer visibility, not only NMP regime.",
          ], title_color=RED_SOFT, size=13.5)
    _card(slide, x0 + 2 * (col_w + gap), Inches(1.5), col_w, col_h, "Known limitations",
          [
              "V(t) is absent from current E(t) values; Dagstuhl has no video. The composite reallocates weight until multimodal data exists.",
              "",
              f"All five E(t) pieces are zero{nbh}latency studio recordings. Cross{nbh}regime variation arrives with Tier{nbh}3.",
              "",
              "p < 0.005 means 0 of 200 permutations exceeded the observed value.",
          ], title_color=GOLD, size=13.5)

    _takeaway(slide, "All retrospective items are documented in sprint3_results.md.")


def slide_8_questions(prs):
    slide = _content_slide(prs, 8, "Problems / questions", "Open questions")
    q_w = SLIDE_W - 2 * MARGIN
    q_h = Inches(1.7)
    items = [
        ("1 · Prof. Hacker: ESMUC dataset access",
         "Do you have institutional access to the ESMUC multitrack dataset? "
         "ChoralSynth is freely available on Zenodo for research and we will download it ourselves; "
         "ESMUC is the only dataset where we need support."),
        ("2 · Coordinators: cluster access (nice-to-have, not a blocker)",
         "Is CPU time available on a Bamberg or HSLU cluster? The planned next-iteration "
         "scope runs overnight on our laptops. Cluster access would let us run denser "
         "jitter grids and finer analysis windows, which strengthens the H1 robustness checks."),
    ]
    y = Inches(1.8)
    for title, body in items:
        _card(slide, MARGIN, y, q_w, q_h, title, [body], size=15)
        y += q_h + Inches(0.35)

    _text(slide, MARGIN, Inches(6.4), q_w, Inches(0.5),
          "Thank you.", size=22, bold=True, color=TEAL,
          align=PP_ALIGN.CENTER)


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_1_title(prs)
    slide_2_recap(prs)
    slide_3_headline(prs)
    slide_4_audio_network(prs)
    slide_5_video_dashboard(prs)
    slide_6_next_iteration(prs)
    slide_7_retro(prs)
    slide_8_questions(prs)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"Wrote {OUT_PATH}")
    print(f"  Slides: {len(prs.slides)}")
    print(f"  Size: {OUT_PATH.stat().st_size / 1024:.1f} KB")


if __name__ == "__main__":
    build()
