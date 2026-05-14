"""Generate apr30_deck.pptx from the deck markdown content.

Run from the project root:
    python scripts/generate_apr30_pptx.py

Writes:
    output/apr30_deck.pptx

The deck content is hardcoded in this script for reproducibility.
Source of truth: apr30_deck.md (8 slides matching Prof. Gloor's Apr 14 email rubric).

Visual style: "Studio Acoustic" palette. Deep teal as the primary, warm gold as
the accent, cream content backgrounds, charcoal body text. Distinct from the
COIN sibling team's navy-and-coral theme; the teal-and-gold direction is
thematically grounded in the project (audio-engineering studio aesthetic,
warm gold echoing the alchemical Citrinitas stage in Gloor's framework).
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# --- Paths ---
ROOT = Path(__file__).resolve().parent.parent
OUT_PATH = ROOT / "output" / "apr30_deck.pptx"

# --- Color palette (Studio Acoustic) ---
# Names are kept short for readability; values define the actual brand.
TEAL = RGBColor(0x0E, 0x4D, 0x5E)          # primary, deep teal
TEAL_DARK = RGBColor(0x07, 0x33, 0x40)     # title-slide background
TEAL_MID = RGBColor(0x32, 0x80, 0x95)      # accents on dark backgrounds
GOLD = RGBColor(0xC4, 0x90, 0x2A)          # warm gold accent (Citrinitas)
GOLD_SOFT = RGBColor(0xE6, 0xC6, 0x70)     # secondary gold tint
CREAM = RGBColor(0xFA, 0xF6, 0xEC)         # content slide background
CREAM_DEEP = RGBColor(0xF1, 0xE9, 0xD3)    # alt row tint
IVORY = RGBColor(0xFD, 0xFB, 0xF5)         # near-white body for cards
CHARCOAL = RGBColor(0x24, 0x26, 0x2E)      # body text
MUTED = RGBColor(0x6E, 0x6A, 0x60)         # warm muted grey for italics
MIST = RGBColor(0xCB, 0xD9, 0xDD)          # subtle teal tint for dark-bg captions

# Status dots
GREEN = RGBColor(0x4F, 0x8E, 0x5A)
AMBER = RGBColor(0xC9, 0x8A, 0x23)
GREY = RGBColor(0xB1, 0xAE, 0xA4)

# --- Backwards-compatible aliases used by the slide-builder code below ---
# (the body of the script was originally written against NAVY/RED/etc; this
#  alias block lets the new palette flow through without renaming every call)
NAVY = TEAL
NAVY_DARK = TEAL_DARK
RED = GOLD
LIGHT_BLUE = MIST
WHITE = IVORY
OFF_WHITE = CREAM

# --- Layout constants (16:9) ---
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.55)
TITLE_TOP = Inches(0.45)
ACCENT_LINE_Y = Inches(1.05)
CONTENT_TOP = Inches(1.25)

# --- Fonts ---
TITLE_FONT = "Calibri"
BODY_FONT = "Calibri"

# ---------- Helpers ----------


def set_slide_bg(slide, color: RGBColor) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def fill_solid(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def no_line(shape) -> None:
    shape.line.fill.background()


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    size: int = 14,
    bold: bool = False,
    color: RGBColor = CHARCOAL,
    font: str = BODY_FONT,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    italic: bool = False,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return tb


def add_rich_paragraphs(
    slide,
    left,
    top,
    width,
    height,
    paragraphs,
    *,
    default_size: int = 14,
    default_color: RGBColor = CHARCOAL,
    line_spacing: float = 1.15,
):
    """Add a textbox with multiple paragraphs.

    Each paragraph is a dict with keys:
      runs: list of dicts {text, size, bold, italic, color, font}
      space_before, space_after: optional Pt values
      bullet: optional bool (adds a bullet character)
      align: optional alignment
    """
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = para.get("align", PP_ALIGN.LEFT)
        if "space_before" in para:
            p.space_before = Pt(para["space_before"])
        if "space_after" in para:
            p.space_after = Pt(para["space_after"])
        p.line_spacing = line_spacing
        bullet_prefix = "•  " if para.get("bullet") else ""
        first = True
        for run_spec in para.get("runs", []):
            run = p.add_run()
            text = run_spec["text"]
            if first and bullet_prefix:
                text = bullet_prefix + text
                first = False
            run.text = text
            run.font.size = Pt(run_spec.get("size", default_size))
            run.font.bold = run_spec.get("bold", False)
            run.font.italic = run_spec.get("italic", False)
            run.font.color.rgb = run_spec.get("color", default_color)
            run.font.name = run_spec.get("font", BODY_FONT)
    return tb


def apply_content_background(slide) -> None:
    """Apply the cream content-slide background. Call on every content slide."""
    set_slide_bg(slide, CREAM)


def add_slide_title(slide, title_text: str, breadcrumb: str | None = None) -> None:
    """Add a slide title with a thin GOLD vertical accent on the left side.

    The vertical-left accent is intentionally distinct from the COIN sibling
    deck's horizontal red bar under the title.
    """
    if breadcrumb:
        add_text(
            slide,
            MARGIN + Inches(0.18),
            Inches(0.35),
            SLIDE_W - 2 * MARGIN - Inches(0.18),
            Inches(0.3),
            breadcrumb,
            size=10,
            color=GOLD,
            bold=True,
        )
        title_top = Inches(0.62)
    else:
        title_top = TITLE_TOP

    add_text(
        slide,
        MARGIN + Inches(0.18),
        title_top,
        SLIDE_W - 2 * MARGIN - Inches(0.18),
        Inches(0.55),
        title_text,
        size=30,
        bold=True,
        color=TEAL,
        font=TITLE_FONT,
    )
    # Vertical gold accent on the left
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        MARGIN,
        Inches(0.32),
        Inches(0.08),
        Inches(0.95),
    )
    fill_solid(accent, GOLD)
    no_line(accent)


def add_status_dot(slide, left, top, status: str):
    color_map = {"green": GREEN, "amber": AMBER, "grey": GREY}
    size = Inches(0.18)
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    fill_solid(shape, color_map.get(status, GREY))
    no_line(shape)
    return shape


def make_table(
    slide,
    left,
    top,
    width,
    height,
    headers,
    rows,
    *,
    col_widths_pct=None,
    header_color: RGBColor = NAVY,
    header_text: RGBColor = WHITE,
    body_size: int = 11,
    header_size: int = 11,
    alt_row_color: RGBColor = OFF_WHITE,
    body_color: RGBColor = CHARCOAL,
    row_height: float | None = None,
):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    if col_widths_pct:
        total = sum(col_widths_pct)
        for i, pct in enumerate(col_widths_pct):
            table.columns[i].width = int(width * pct / total)

    if row_height is not None:
        for r in range(n_rows):
            table.rows[r].height = Inches(row_height)

    # Header row
    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        cell.margin_left = Inches(0.1)
        cell.margin_right = Inches(0.1)
        cell.margin_top = Inches(0.04)
        cell.margin_bottom = Inches(0.04)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.text = ""
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = header
        run.font.size = Pt(header_size)
        run.font.bold = True
        run.font.color.rgb = header_text
        run.font.name = TITLE_FONT

    # Body rows
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = OFF_WHITE if r % 2 == 0 else WHITE
            cell.margin_left = Inches(0.1)
            cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.05)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.text = ""
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT

            # Support 2-tuple (text, bold) for first-cell emphasis
            if isinstance(val, tuple):
                text_val, is_bold = val[0], val[1] if len(val) > 1 else False
            else:
                text_val = str(val)
                is_bold = False

            run = p.add_run()
            run.text = text_val
            run.font.size = Pt(body_size)
            run.font.bold = is_bold
            run.font.color.rgb = body_color
            run.font.name = BODY_FONT

    return table_shape


# ---------- Build the deck ----------


def build_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    build_slide_1_title(prs.slides.add_slide(blank))
    build_slide_2_team(prs.slides.add_slide(blank))
    build_slide_3_goals(prs.slides.add_slide(blank))
    build_slide_4_dataset(prs.slides.add_slide(blank))
    build_slide_5_plan(prs.slides.add_slide(blank))
    build_slide_6_iteration(prs.slides.add_slide(blank))
    build_slide_7_wow(prs.slides.add_slide(blank))
    build_slide_8_thanks(prs.slides.add_slide(blank))

    return prs


# ---------- Slide 1: Title ----------


def build_slide_1_title(slide) -> None:
    set_slide_bg(slide, NAVY_DARK)

    # Subtle red accent bar
    acc = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.7),
        Inches(2.65),
        Inches(0.7),
        Inches(0.07),
    )
    fill_solid(acc, RED)
    no_line(acc)

    add_text(
        slide,
        Inches(0.7),
        Inches(2.85),
        Inches(11.9),
        Inches(1.4),
        "Status Meeting II",
        size=64,
        bold=True,
        color=RED,
        font="Calibri",
    )

    add_text(
        slide,
        Inches(0.7),
        Inches(4.3),
        Inches(11.9),
        Inches(0.55),
        "Project 8: Entanglement in Online Choir",
        size=22,
        color=WHITE,
    )

    add_text(
        slide,
        Inches(0.7),
        Inches(5.0),
        Inches(11.9),
        Inches(0.4),
        "SNA-OSN-M Summer 2026 · Uni Bamberg × Uni Köln × HSLU",
        size=13,
        color=LIGHT_BLUE,
    )

    add_text(
        slide,
        Inches(0.7),
        Inches(5.5),
        Inches(11.9),
        Inches(0.4),
        "Presented by Zuraiz, on behalf of the team",
        size=13,
        color=LIGHT_BLUE,
    )

    add_text(
        slide,
        Inches(0.7),
        Inches(5.95),
        Inches(11.9),
        Inches(0.4),
        "Supervisors: Prof. Janine Hacker (Uni Bamberg) · Prof. Peter Gloor (MIT, COIN seminar lead)",
        size=12,
        color=LIGHT_BLUE,
    )

    add_text(
        slide,
        Inches(0.7),
        Inches(6.85),
        Inches(11.9),
        Inches(0.35),
        "2026-04-30 · 14:00 CET",
        size=11,
        color=LIGHT_BLUE,
        italic=True,
    )


# ---------- Slide 2: Team ----------


def build_slide_2_team(slide) -> None:
    apply_content_background(slide)
    add_slide_title(slide, "The Team", breadcrumb="INTRODUCTION")

    members = [
        ("Zuraiz", "Tree Hugger 73% Ant", "Ant 38% / Capybara 38%"),
        ("Hammad Anwar", "Tree Hugger 64% Ant", "Capybara 49%"),
        ("Hassan Ahmed", "Tree Hugger 63% Ant", "Ant 40%"),
        ("Kumaran Vasu", "Tree Hugger 75% Ant", "Ant 80%"),
    ]

    n = len(members)
    band_top = Inches(1.55)
    portrait_size = Inches(1.7)

    # Total horizontal space for portraits + gaps
    total_w = SLIDE_W - 2 * MARGIN
    per_col = total_w / n
    for i, (name, sc_chat, symbiont) in enumerate(members):
        col_left = MARGIN + i * per_col
        # Portrait circle (placeholder)
        circle_left = col_left + (per_col - portrait_size) / 2
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            int(circle_left),
            int(band_top),
            portrait_size,
            portrait_size,
        )
        fill_solid(circle, LIGHT_BLUE)
        circle.line.color.rgb = NAVY
        circle.line.width = Pt(1.5)

        # Initials in the circle
        initials = "".join(part[0] for part in name.split()[:2]).upper()
        add_text(
            slide,
            int(circle_left),
            int(band_top),
            portrait_size,
            portrait_size,
            initials,
            size=44,
            bold=True,
            color=NAVY,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
            font=TITLE_FONT,
        )

        # Name below circle
        name_top = band_top + portrait_size + Inches(0.15)
        add_text(
            slide,
            int(col_left),
            int(name_top),
            int(per_col),
            Inches(0.4),
            name,
            size=15,
            bold=True,
            color=NAVY,
            align=PP_ALIGN.CENTER,
        )

        # SC Chat label/value
        sc_label_top = name_top + Inches(0.5)
        add_text(
            slide,
            int(col_left + Inches(0.15)),
            int(sc_label_top),
            int(per_col - Inches(0.3)),
            Inches(0.3),
            f"SC Chat:  {sc_chat}",
            size=11,
            color=CHARCOAL,
            align=PP_ALIGN.CENTER,
        )

        # Symbiont label/value
        sb_label_top = sc_label_top + Inches(0.35)
        add_text(
            slide,
            int(col_left + Inches(0.15)),
            int(sb_label_top),
            int(per_col - Inches(0.3)),
            Inches(0.3),
            f"Symbiont:  {symbiont}",
            size=11,
            color=CHARCOAL,
            align=PP_ALIGN.CENTER,
        )

    # Footer text
    add_rich_paragraphs(
        slide,
        MARGIN,
        Inches(5.55),
        SLIDE_W - 2 * MARGIN,
        Inches(1.6),
        [
            {
                "runs": [
                    {
                        "text": "All four members are M.Sc. students at Uni Bamberg. ",
                        "size": 13,
                    },
                    {
                        "text": "Three carried over from the chapter team; Kumaran joined at the project signup. ",
                        "size": 13,
                    },
                ]
            },
            {
                "space_before": 8,
                "runs": [
                    {
                        "text": "The archetype results are shown because the May 21 Virtual Mirror requires us to analyse our own team communication, not because they are the scientific contribution today.",
                        "size": 13,
                        "italic": True,
                        "color": MUTED,
                    },
                ],
            },
        ],
    )


# ---------- Slide 3: Goals (Scope) ----------


def build_slide_3_goals(slide) -> None:
    apply_content_background(slide)
    add_slide_title(slide, "Goals (Scope)", breadcrumb="WHY THIS PROJECT")

    # Research question banner
    rq_top = Inches(1.45)
    rq_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        MARGIN,
        rq_top,
        SLIDE_W - 2 * MARGIN,
        Inches(0.85),
    )
    fill_solid(rq_bg, NAVY)
    no_line(rq_bg)

    add_text(
        slide,
        MARGIN + Inches(0.25),
        rq_top + Inches(0.1),
        SLIDE_W - 2 * MARGIN - Inches(0.5),
        Inches(0.3),
        "RESEARCH QUESTION",
        size=10,
        bold=True,
        color=LIGHT_BLUE,
    )
    add_text(
        slide,
        MARGIN + Inches(0.25),
        rq_top + Inches(0.4),
        SLIDE_W - 2 * MARGIN - Inches(0.5),
        Inches(0.5),
        "When a choir sings together over the internet, can we put a number on how well coordinated they are?",
        size=18,
        bold=True,
        color=WHITE,
    )

    # Two-column layout below
    col_top = Inches(2.55)
    col_height = Inches(4.2)
    gutter = Inches(0.4)
    col_w = (SLIDE_W - 2 * MARGIN - gutter) / 2

    # Left column: Where this comes from
    left_x = MARGIN
    add_text(
        slide,
        left_x,
        col_top,
        col_w,
        Inches(0.35),
        "WHERE THIS COMES FROM",
        size=12,
        bold=True,
        color=NAVY,
    )

    add_rich_paragraphs(
        slide,
        left_x,
        col_top + Inches(0.4),
        col_w,
        col_height - Inches(0.4),
        [
            {
                "runs": [
                    {"text": "Research question: ", "bold": True, "size": 13},
                    {
                        "text": "from the gap between NMP studies (latency, tool quality) and coordination research (group flow). Choirs give an unusually concrete acoustic outcome, measurable in milliseconds.",
                        "size": 13,
                    },
                ],
                "space_after": 14,
            },
            {
                "runs": [
                    {"text": "E(t) formula: ", "bold": True, "size": 13},
                    {
                        "text": "inspired by Gloor's entanglement work on team communication and Pentland's Honest Signals, ",
                        "size": 13,
                    },
                    {
                        "text": "but adapted by us for choir audio and video. ",
                        "size": 13,
                        "bold": True,
                    },
                    {
                        "text": "The original entanglement formula was validated on email rhythms, not music, so H1 to H3 are our validation tests.",
                        "size": 13,
                    },
                ],
            },
        ],
    )

    # Right column: Three hypotheses
    right_x = MARGIN + col_w + gutter
    add_text(
        slide,
        right_x,
        col_top,
        col_w,
        Inches(0.35),
        "THREE FALSIFIABLE HYPOTHESES",
        size=12,
        bold=True,
        color=NAVY,
    )

    hyps = [
        (
            "H1",
            "Low-latency tools (Jamulus, SoundJack) score higher on E(t) than high-latency tools (Zoom).",
        ),
        (
            "H2",
            'The "who influences whom" network shifts shape between low and high latency, from democratic to leader-dominated.',
        ),
        (
            "H3",
            "Adding body-movement signals on top of audio improves how E(t) tracks coordination, by at least 10 percentage points of explained variance.",
        ),
    ]

    h_y = col_top + Inches(0.4)
    for label, body in hyps:
        # Label badge
        badge = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            right_x,
            h_y,
            Inches(0.55),
            Inches(0.45),
        )
        fill_solid(badge, RED)
        no_line(badge)
        add_text(
            slide,
            right_x,
            h_y,
            Inches(0.55),
            Inches(0.45),
            label,
            size=14,
            bold=True,
            color=WHITE,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )

        # Body text next to badge
        add_text(
            slide,
            right_x + Inches(0.7),
            h_y,
            col_w - Inches(0.7),
            Inches(0.95),
            body,
            size=12,
            color=CHARCOAL,
            anchor=MSO_ANCHOR.TOP,
        )
        h_y += Inches(1.05)

    # Bottom strip
    bottom_y = Inches(6.95)
    add_text(
        slide,
        MARGIN,
        bottom_y,
        SLIDE_W - 2 * MARGIN,
        Inches(0.4),
        "Either way the result is publishable. Null finding: latency dichotomy is too crude. Positive finding: the field gets its first quantitative coordination meter.",
        size=11,
        italic=True,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )


# ---------- Slide 4: Dataset Strategy ----------


def build_slide_4_dataset(slide) -> None:
    apply_content_background(slide)
    add_slide_title(slide, "Dataset Strategy", breadcrumb="WHAT DATA WE USE")

    # Mental model strap
    mm_y = Inches(1.4)
    add_text(
        slide,
        MARGIN,
        mm_y,
        SLIDE_W - 2 * MARGIN,
        Inches(0.4),
        "Mental model:  Tier 1 is real-world but messy.  Tier 2 is clean but academic.  Tier 3 is controlled and experimental.",
        size=13,
        italic=True,
        color=NAVY,
        bold=True,
    )

    # Tier table
    table_top = Inches(1.85)
    table_w = SLIDE_W - 2 * MARGIN

    headers = ["Tier", "What it is", "What we do with it", "What it can prove"]
    rows = [
        (
            ("Tier 0\nseed URLs", True),
            "5 Jamulus or Choir@Home URLs from Prof. Hacker, already in hand",
            "Examples that guide search terms and inclusion rules",
            "Nothing alone; starts the corpus search",
        ),
        (
            ("Tier 1\nYouTube virtual choirs", True),
            "20 to 30 public videos, hand-curated by May 15",
            "Body movement, mouth opening, visible synchrony, ensemble audio",
            "Visual coordination and H3 sanity check",
        ),
        (
            ("Tier 2\nacademic multitrack", True),
            "Dagstuhl ChoirSet, ESMUC, ChoralSynth",
            "Separate singer tracks for pitch, onset timing, Granger influence",
            "Per-singer audio: A(t), N(t), influence graph",
        ),
        (
            ("Tier 3\nlatency injection", True),
            "Synthesised Zoom-class and low-latency variants of Tier 2",
            "Add delay, jitter, packet loss, recompute E(t)",
            "Clean H1 test, injected latency is known",
        ),
    ]

    make_table(
        slide,
        MARGIN,
        table_top,
        table_w,
        Inches(3.7),
        headers,
        rows,
        col_widths_pct=[16, 28, 28, 28],
        body_size=11,
        header_size=12,
        row_height=0.74,
    )

    # Limitation note (positioned with safety buffer below the table)
    note_y = Inches(5.95)
    add_rich_paragraphs(
        slide,
        MARGIN,
        note_y,
        SLIDE_W - 2 * MARGIN,
        Inches(1.4),
        [
            {
                "runs": [
                    {"text": "Important limitation:  ", "bold": True, "size": 12, "color": RED},
                    {
                        "text": "Tier 1 YouTube audio is mixed stereo, so it cannot support per-singer Granger networks. That is why H1 and the influence graph rely on Tier 2 and Tier 3.",
                        "size": 12,
                    },
                ],
                "space_after": 6,
            },
            {
                "runs": [
                    {"text": "Guardrail:  ", "bold": True, "size": 12, "color": NAVY},
                    {
                        "text": "no self-recording, no third-party choir recruitment in this project phase. Public videos and published datasets only; only derived features kept after extraction. GDPR DPIA outline due to Bamberg DPO by May 21.",
                        "size": 12,
                    },
                ],
            },
        ],
    )


# ---------- Slide 5: Overall Project Plan ----------


def build_slide_5_plan(slide) -> None:
    apply_content_background(slide)
    add_slide_title(slide, "Overall Project Plan", breadcrumb="14 WEEKS, 3 PHASES")

    # Phase strap (3 horizontal bars)
    strap_top = Inches(1.4)
    strap_h = Inches(0.5)
    strap_w_total = SLIDE_W - 2 * MARGIN
    phase_widths = [
        ("Phase 1: Scope and Scaffold", "Apr 16 - Apr 30", NAVY, 0.16),
        ("Phase 2: Build and Analyse", "May 1 - Jul 7", RGBColor(0x36, 0x5C, 0xA8), 0.55),
        ("Phase 3: Synthesise", "Jul 8 - Jul 31", RGBColor(0x6F, 0x8F, 0xC8), 0.29),
    ]
    x = MARGIN
    for label, dates, color, frac in phase_widths:
        w = int(strap_w_total * frac)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, strap_top, w, strap_h)
        fill_solid(bar, color)
        no_line(bar)
        add_text(
            slide,
            x + Inches(0.15),
            strap_top + Inches(0.06),
            w - Inches(0.3),
            Inches(0.22),
            label,
            size=11,
            bold=True,
            color=WHITE,
        )
        add_text(
            slide,
            x + Inches(0.15),
            strap_top + Inches(0.27),
            w - Inches(0.3),
            Inches(0.2),
            dates,
            size=10,
            color=LIGHT_BLUE,
        )
        x += w

    # Status meeting calendar
    cal_top = Inches(2.05)
    add_text(
        slide,
        MARGIN,
        cal_top,
        SLIDE_W - 2 * MARGIN,
        Inches(0.3),
        "STATUS MEETINGS",
        size=10,
        bold=True,
        color=MUTED,
    )

    # Calendar markers
    cal_band_top = Inches(2.4)
    cal_w = SLIDE_W - 2 * MARGIN
    markers = [
        ("VS#1", "Apr 16", "Kick-off", False),
        ("VS#2", "Apr 30", "Today", True),
        ("VS#3", "May 21", "Mirror", False),
        ("VS#4", "Jun 11", "", False),
        ("VS#5", "Jun 25", "", False),
        ("VS#6", "Jul 9", "Pre-final", False),
        ("Final", "Jul 23", "Demo", False),
        ("Paper", "Jul 31", "Due", False),
    ]
    n = len(markers)
    per_w = cal_w / n
    for i, (label, date, note, is_today) in enumerate(markers):
        cx = MARGIN + i * per_w + per_w / 2
        circle_size = Inches(0.55)
        circle_left = int(cx - circle_size / 2)
        circle_top = cal_band_top
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            circle_left,
            circle_top,
            circle_size,
            circle_size,
        )
        if is_today:
            fill_solid(circle, RED)
            circle.line.color.rgb = RED
        else:
            fill_solid(circle, NAVY)
            circle.line.color.rgb = NAVY
        circle.line.width = Pt(0)

        add_text(
            slide,
            circle_left,
            circle_top,
            circle_size,
            circle_size,
            label.replace("VS#", "#"),
            size=10,
            bold=True,
            color=WHITE,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )

        add_text(
            slide,
            int(cx - per_w / 2),
            circle_top + circle_size + Inches(0.06),
            per_w,
            Inches(0.25),
            date,
            size=10,
            bold=True,
            color=NAVY if not is_today else RED,
            align=PP_ALIGN.CENTER,
        )

        if note:
            add_text(
                slide,
                int(cx - per_w / 2),
                circle_top + circle_size + Inches(0.3),
                per_w,
                Inches(0.25),
                note,
                size=9,
                italic=True,
                color=MUTED,
                align=PP_ALIGN.CENTER,
            )

    # Milestones
    mile_title_top = Inches(3.55)
    add_text(
        slide,
        MARGIN,
        mile_title_top,
        SLIDE_W - 2 * MARGIN,
        Inches(0.3),
        "PHASE-BASED MILESTONES",
        size=10,
        bold=True,
        color=MUTED,
    )

    mile_top = Inches(3.85)
    mile_height = Inches(3.45)
    mile_w = SLIDE_W - 2 * MARGIN

    milestones = [
        ("Research question, hypotheses, and data tiers defined", "1", "green"),
        ("Repo scaffold ready (package, CI, smoke tests)", "1", "green"),
        ("Current data state recorded: seed URLs only", "1", "green"),
        ("Tier 2 academic datasets on disk with manifest", "2", "amber"),
        ("Tier 1 YouTube corpus curated with inclusion log", "2", "grey"),
        ("Video pipeline producing per-singer pose data", "2", "grey"),
        ("Network pipeline producing influence graph", "2", "grey"),
        ("E(t) computed end-to-end on full corpus", "2", "grey"),
        ("Dashboard alpha (60-second live demo target)", "2", "grey"),
        ("Paper draft v1", "3", "grey"),
        ("Final presentation and final paper submitted", "3", "grey"),
    ]

    # Build a custom table for milestones with status dots in a column
    headers = ["Milestone", "Phase", "Status"]
    rows = [(m, p, "") for m, p, _ in milestones]

    table_shape = make_table(
        slide,
        MARGIN,
        mile_top,
        mile_w,
        mile_height,
        headers,
        rows,
        col_widths_pct=[78, 11, 11],
        body_size=11,
        header_size=11,
        row_height=0.27,
    )

    # Overlay status dots into the third column
    table = table_shape.table
    # Compute approximate positions for the status dots
    third_col_x = MARGIN + int(mile_w * 0.89) + Inches(0.05)  # after 78+11=89%
    third_col_center_x = MARGIN + int(mile_w * 0.945)  # midpoint of 89-100
    row_height_emu = Emu(int(Inches(0.27)))
    header_h = Inches(0.27)
    for i, (_, _, status) in enumerate(milestones):
        row_top = mile_top + header_h + i * Inches(0.27)
        dot_left = third_col_center_x - Inches(0.09)
        dot_top = row_top + Inches(0.05)
        add_status_dot(slide, dot_left, dot_top, status)


# ---------- Slide 6: Plan for the Next Iteration ----------


def build_slide_6_iteration(slide) -> None:
    apply_content_background(slide)
    add_slide_title(slide, "Plan for the Next Iteration", breadcrumb="SPRINT 2: APR 30 TO MAY 21")

    # Sprint header callout
    sh_top = Inches(1.4)
    sh_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        MARGIN,
        sh_top,
        SLIDE_W - 2 * MARGIN,
        Inches(0.55),
    )
    fill_solid(sh_bg, OFF_WHITE)
    sh_bg.line.color.rgb = NAVY
    sh_bg.line.width = Pt(0.5)

    add_text(
        slide,
        MARGIN + Inches(0.2),
        sh_top + Inches(0.12),
        SLIDE_W - 2 * MARGIN - Inches(0.4),
        Inches(0.35),
        "3 weeks · ends at VS#3 Mirror Session on May 21 · 4 work packages, each ending in a defined deliverable",
        size=12,
        bold=True,
        color=NAVY,
    )

    # Sprint table (compact body to leave clean room for the bottom note)
    table_top = Inches(2.1)

    headers = ["Work Package", "Next Step", "Output"]
    rows = [
        (
            ("Audio pipeline", True),
            "Pull academic multitrack datasets (Dagstuhl, ESMUC, ChoralSynth). Run pitch tracking and onset detection.",
            "Per-singer feature parquet for 5 pieces, by May 8.",
        ),
        (
            ("Video pipeline", True),
            "Run MediaPipe Pose and FaceMesh on the first 10 YouTube videos. Calibrate against a published pose ground truth.",
            "Per-singer pose parquet plus a 1-page calibration note, by May 22.",
        ),
        (
            ("Curated corpus + GDPR", True),
            "Hand-curate 20 to 30 YouTube virtual-choir videos. File the GDPR DPIA outline (face-mesh landmarks are biometric data under Art. 9).",
            "Corpus manifest with SHA-256 hashes by May 15. DPIA outline by May 21.",
        ),
        (
            ("Virtual Mirror (seminar)", True),
            "Export the team WhatsApp chat. Run SocialCompass tools (SC Chat, Symbiont, Beecome). Classify the team archetype with all 4 members.",
            "Mirror write-up with archetype classification, presented at VS#3 on May 21.",
        ),
    ]

    make_table(
        slide,
        MARGIN,
        table_top,
        SLIDE_W - 2 * MARGIN,
        Inches(4.45),
        headers,
        rows,
        col_widths_pct=[20, 50, 30],
        body_size=10,
        header_size=12,
        row_height=0.92,
    )

    # Bottom note (placed with generous buffer; 6.85 leaves ~0.4" before slide bottom)
    add_text(
        slide,
        MARGIN,
        Inches(6.85),
        SLIDE_W - 2 * MARGIN,
        Inches(0.45),
        "Dataset construction order in Sprint 2:  Tier 2 on disk first, Tier 1 curated second, Tier 3 generated from Tier 2 once the audio pipeline can run.",
        size=10,
        italic=True,
        color=MUTED,
    )


# ---------- Slide 7: Way of Working ----------


def build_slide_7_wow(slide) -> None:
    apply_content_background(slide)
    add_slide_title(slide, "Way of Working", breadcrumb="HOW THE TEAM OPERATES")

    columns = [
        (
            "Cadence",
            [
                ("Format", "3-week sprints aligned to status meetings"),
                ("Sync", "Pre-VS review 1 to 2 days before each meeting"),
                ("Delivery", "Each sprint ends in a defined artefact"),
                ("Meeting", "VS with supervisors per sprint"),
            ],
        ),
        (
            "Sync",
            [
                ("Weekly", "30-minute team sync"),
                ("Daily", "Async status post per person"),
                ("Channels", "WhatsApp group + GitHub Issues"),
                ("Supervisor", "Weekly check-in with Prof. Hacker"),
            ],
        ),
        (
            "Toolstack",
            [
                ("Code + CI", "Python, Git, GitHub Actions"),
                ("Docs", "Markdown source of truth + Obsidian method notes"),
                ("Quality", "Pre-commit hooks (ruff, mypy), peer review via PR"),
                ("Meetings", "Zoom for VS, screen-share for live demos"),
            ],
        ),
    ]

    n = len(columns)
    col_top = Inches(1.45)
    col_height = Inches(4.85)
    gutter = Inches(0.3)
    col_w = (SLIDE_W - 2 * MARGIN - (n - 1) * gutter) / n

    for i, (header, items) in enumerate(columns):
        x = MARGIN + i * (col_w + gutter)

        # Card background (ivory) with a thick TEAL vertical stripe on the left.
        # This is intentionally distinct from the COIN sibling deck's
        # filled-navy header bar over light-blue body.
        card = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x,
            col_top,
            col_w,
            col_height,
        )
        fill_solid(card, IVORY)
        card.line.color.rgb = GOLD_SOFT
        card.line.width = Pt(0.5)

        stripe = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x,
            col_top,
            Inches(0.12),
            col_height,
        )
        fill_solid(stripe, TEAL)
        no_line(stripe)

        # Card header (text-only, no filled bar)
        add_text(
            slide,
            x + Inches(0.35),
            col_top + Inches(0.25),
            col_w - Inches(0.5),
            Inches(0.5),
            header,
            size=20,
            bold=True,
            color=TEAL,
        )
        # Thin gold underline beneath header
        underline = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x + Inches(0.35),
            col_top + Inches(0.78),
            Inches(0.6),
            Inches(0.03),
        )
        fill_solid(underline, GOLD)
        no_line(underline)

        # Items below header
        item_y = col_top + Inches(0.95)
        for label, val in items:
            add_text(
                slide,
                x + Inches(0.35),
                item_y,
                col_w - Inches(0.5),
                Inches(0.3),
                label,
                size=11,
                bold=True,
                color=GOLD,
            )
            add_text(
                slide,
                x + Inches(0.35),
                item_y + Inches(0.3),
                col_w - Inches(0.5),
                Inches(0.55),
                val,
                size=12,
                color=CHARCOAL,
            )
            item_y += Inches(0.95)

    # Bottom note (full width)
    add_text(
        slide,
        MARGIN,
        Inches(6.55),
        SLIDE_W - 2 * MARGIN,
        Inches(0.7),
        "The team WhatsApp group, set up per Prof. Gloor's Apr 14 instruction, is also our input data for the Virtual Mirror on May 21. The way we communicate becomes the data we analyse on ourselves.",
        size=11,
        italic=True,
        color=MUTED,
    )


# ---------- Slide 8: Thank you ----------


def build_slide_8_thanks(slide) -> None:
    set_slide_bg(slide, NAVY_DARK)

    # Accent
    acc = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.7),
        Inches(2.85),
        Inches(0.7),
        Inches(0.07),
    )
    fill_solid(acc, RED)
    no_line(acc)

    add_text(
        slide,
        Inches(0.7),
        Inches(3.05),
        Inches(11.9),
        Inches(1.4),
        "Thank you for listening.",
        size=58,
        bold=True,
        color=RED,
    )

    add_text(
        slide,
        Inches(0.7),
        Inches(4.55),
        Inches(11.9),
        Inches(0.6),
        "Open to your questions.",
        size=24,
        color=WHITE,
    )

    add_text(
        slide,
        Inches(0.7),
        Inches(6.55),
        Inches(11.9),
        Inches(0.4),
        "Project 8 · Entanglement in Online Choir · Zuraiz · m.zuraiz2001@gmail.com",
        size=11,
        color=LIGHT_BLUE,
        italic=True,
    )


# ---------- Main ----------


def main() -> None:
    prs = build_deck()
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))
    print(f"Wrote {OUT_PATH}")


if __name__ == "__main__":
    main()
