"""Generate the final Jul-23 presentation from the audited deck source."""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

ROOT = Path(__file__).resolve().parent.parent
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from scripts.generate_jun11_pptx import (  # noqa: E402
    CREAM,
    FIG_DIR,
    GOLD_SOFT,
    IVORY,
    MARGIN,
    MIST,
    MUTED,
    SLIDE_H,
    SLIDE_W,
    TEAL,
    TEAL_DARK,
    _card,
    _kicker,
    _picture_fit,
    _solid_bg,
    _stat_card,
    _takeaway,
    _text,
    _title,
)

OUT_PATH = Path("output/jul23_final_presentation.pptx")
TOTAL_SLIDES = 19


def _content_slide(prs: Presentation, n: int, kicker: str, title: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, CREAM)
    _kicker(slide, kicker)
    _title(slide, title)
    _text(
        slide,
        SLIDE_W - Inches(1.7),
        SLIDE_H - Inches(0.38),
        Inches(1.2),
        Inches(0.3),
        f"{n} / {TOTAL_SLIDES}",
        size=10,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )
    _text(
        slide,
        MARGIN,
        SLIDE_H - Inches(0.38),
        Inches(5.0),
        Inches(0.3),
        "Project 8 - Entanglement in Online Choir",
        size=10,
        color=MUTED,
    )
    return slide


def _stack(slide, rows, *, y=1.5, height=0.92, gap=0.15, size=13.5):
    current_y = Inches(y)
    for heading, body in rows:
        _card(
            slide,
            MARGIN,
            current_y,
            SLIDE_W - 2 * MARGIN,
            Inches(height),
            heading,
            [body],
            size=size,
        )
        current_y += Inches(height + gap)


def s1_title(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, TEAL_DARK)
    _text(
        slide,
        MARGIN,
        Inches(1.55),
        SLIDE_W - 2 * MARGIN,
        Inches(1.2),
        "Measuring Coordination in Online Choirs",
        size=44,
        bold=True,
        color=IVORY,
        align=PP_ALIGN.CENTER,
    )
    _text(
        slide,
        MARGIN,
        Inches(2.95),
        SLIDE_W - 2 * MARGIN,
        Inches(0.55),
        "Project 8: Entanglement in Online Choir",
        size=25,
        color=GOLD_SOFT,
        align=PP_ALIGN.CENTER,
    )
    _text(
        slide,
        MARGIN,
        Inches(4.15),
        SLIDE_W - 2 * MARGIN,
        Inches(0.4),
        "Zuraiz - Hammad Anwar - Hassan Ahmed - Kumaran Vasu",
        size=16,
        bold=True,
        color=IVORY,
        align=PP_ALIGN.CENTER,
    )
    _text(
        slide,
        MARGIN,
        Inches(4.85),
        SLIDE_W - 2 * MARGIN,
        Inches(0.4),
        "Supervisors: Prof. Janine Hacker - Prof. Peter Gloor",
        size=14,
        color=MIST,
        align=PP_ALIGN.CENTER,
    )
    _text(
        slide,
        MARGIN,
        Inches(5.65),
        SLIDE_W - 2 * MARGIN,
        Inches(0.4),
        "2026-07-23",
        size=14,
        color=GOLD_SOFT,
        align=PP_ALIGN.CENTER,
    )


def s2_goals(prs: Presentation) -> None:
    slide = _content_slide(prs, 2, "Goals", "What we set out to test")
    _text(
        slide,
        MARGIN,
        Inches(1.35),
        SLIDE_W - 2 * MARGIN,
        Inches(0.55),
        "Measure network effects on choir coordination and deliver an auditable analysis system.",
        size=18,
        bold=True,
        color=TEAL,
    )
    _stack(
        slide,
        [
            (
                "H1 - latency",
                "Higher latency reduces coordination; primary outcome: zero-lag onset synchrony.",
            ),
            (
                "H2 - leadership",
                "Influence networks show leadership structure; outcome: out-degree Gini versus matched random graphs.",
            ),
            (
                "H3 - visual information",
                "Visual signals add information beyond audio; the full test compares audio-only with audio-plus-visual prediction.",
            ),
        ],
        y=2.15,
        height=1.0,
    )
    _takeaway(
        slide,
        "Directions were declared at project start; metric revisions are dated and documented.",
    )


def s3_related(prs: Presentation) -> None:
    slide = _content_slide(prs, 3, "Related work", "Three foundations, one research gap")
    _stack(
        slide,
        [
            (
                "Entanglement",
                "Temporal alignment has been measured in organizational communication networks.",
            ),
            (
                "Honest Signals",
                "Activity, consistency, influence, and mimicry can reveal group coordination.",
            ),
            (
                "Choir and NMP research",
                "Multitrack corpora and latency studies provide audio evidence, but no integrated acoustic, visual, and influence measure.",
            ),
            (
                "Research gap",
                "Adapt these ideas to singing, test them against controlled degradation, and state where the data cannot support a claim.",
            ),
        ],
        y=1.5,
        height=1.0,
    )


def s4_structure(prs: Presentation) -> None:
    slide = _content_slide(prs, 4, "Project", "Whole-project structure")
    _picture_fit(
        slide,
        FIG_DIR / "project_structure.png",
        MARGIN,
        Inches(1.35),
        SLIDE_W - 2 * MARGIN,
        Inches(5.35),
    )
    _takeaway(slide, "Every result has a visible path from source data to tested claim.")


def s5_data(prs: Presentation) -> None:
    slide = _content_slide(prs, 5, "Data", "Three tiers and one multimodal constraint")
    _stack(
        slide,
        [
            (
                "Tier 1 - video",
                "29 virtual-choir videos; 18 pose-usable. Mixed audio plus visual features for exploratory H3.",
            ),
            (
                "Tier 2 - multitrack",
                "Dagstuhl 5, ESMUC 3, ChoralSynth 20. Per-singer audio and influence networks.",
            ),
            (
                "Tier 3 - controlled degradation",
                "Five regimes applied to all 28 Tier-2 pieces for the H1 latency experiment.",
            ),
        ],
        y=1.5,
        height=1.0,
    )
    _text(
        slide,
        MARGIN,
        Inches(5.05),
        SLIDE_W - 2 * MARGIN,
        Inches(0.95),
        "No item has synchronized per-singer audio and per-singer video. Tier 1 has mixed audio; Tier 2 has no video.",
        size=16,
        bold=True,
        color=TEAL,
    )
    _takeaway(slide, "Each dataset is used only for the claims it can support.")


def s6_method(prs: Presentation) -> None:
    slide = _content_slide(prs, 6, "Method", "E(t) and the controlled latency grid")
    _text(
        slide,
        MARGIN,
        Inches(1.35),
        SLIDE_W - 2 * MARGIN,
        Inches(0.45),
        "E(t) combines available acoustic A(t), visual V(t), and network N(t) coordination.",
        size=16,
        bold=True,
        color=TEAL,
    )
    regimes = [
        ("Clean", "0 / 0 / 0%"),
        ("In-person", "25 / 10 / 0%"),
        ("Jamulus LAN", "47 / 46 / 1%"),
        ("Jamulus WAN", "83 / 57 / 3%"),
        ("Zoom-class", "150 / 80 / 8%"),
    ]
    x = MARGIN
    for name, values in regimes:
        _stat_card(slide, x, Inches(2.05), Inches(2.32), Inches(1.35), name, values, number_size=16)
        x += Inches(2.44)
    _text(
        slide,
        MARGIN,
        Inches(3.55),
        Inches(5),
        Inches(0.3),
        "Delay ms / jitter SD ms / dropout",
        size=12,
        color=MUTED,
    )
    _card(
        slide,
        MARGIN,
        Inches(4.1),
        SLIDE_W - 2 * MARGIN,
        Inches(1.65),
        "Inference",
        [
            "Envelope and network analyses use 2,000 circular shifts per cell. H1 onset inference uses a paired clean-to-Zoom sign test across pieces."
        ],
        size=14,
    )
    _takeaway(slide, "The same piece is evaluated in every regime and serves as its own control.")


def s7_process(prs: Presentation) -> None:
    slide = _content_slide(prs, 7, "Work process", "How the team organized the project")
    _stack(
        slide,
        [
            (
                "Four work packages",
                "Acoustic/integration, visual extraction, influence networks, and dashboard/report integration.",
            ),
            (
                "Iteration outputs",
                "Every iteration ended in a schema, result table, test, figure, or presentation artifact.",
            ),
            (
                "Review gates",
                "Seeded scripts, committed summaries, and independent audits connected implementation to claims.",
            ),
            (
                "Virtual Mirror",
                "High shared meaning, low emotional exchange, medium relationship intensity; a weekly asynchronous check-in was adopted.",
            ),
        ],
        y=1.45,
        height=1.0,
    )


def s8_repro(prs: Presentation) -> None:
    slide = _content_slide(prs, 8, "Method", "Reproducibility is part of the result")
    _stack(
        slide,
        [
            (
                "One command",
                "`make reproduce` regenerates report-stage statistics, figures, and the PDF.",
            ),
            (
                "50 tests",
                "Scientific pipelines, dashboard consistency, figure definitions, H3 controls, and report generation are covered.",
            ),
            ("Complete grid", "140 final cells from the committed 2,000-shuffle SLURM protocol."),
            (
                "Claim provenance",
                "Every numerical statement in the deck traces to a committed CSV.",
            ),
        ],
        y=1.5,
        height=1.0,
    )
    _takeaway(slide, "The results can be audited without relying on the presentation.")


def s9_h1(prs: Presentation) -> None:
    slide = _content_slide(prs, 9, "Result", "H1: latency breaks note timing")
    stats = [
        ("-56.5%", "Dagstuhl"),
        ("-65.1%", "ESMUC"),
        ("-75.1%", "ChoralSynth"),
        ("-70.7%", "Corpus; 28/28 decrease"),
    ]
    x = MARGIN
    for number, label in stats:
        _stat_card(slide, x, Inches(1.35), Inches(2.9), Inches(1.25), number, label)
        x += Inches(3.05)
    _picture_fit(
        slide, FIG_DIR / "tier3_corpus_summary.png", MARGIN, Inches(2.85), Inches(7.5), Inches(3.55)
    )
    _card(
        slide,
        Inches(8.25),
        Inches(2.95),
        Inches(4.45),
        Inches(3.25),
        "Channel comparison",
        [
            "Pure envelope A(t): -7.9% Dagstuhl, -11.9% corpus. The envelope-plus-network E(t) composite can rise because network density offsets the envelope decline."
        ],
        size=13.5,
    )
    _takeaway(slide, "All 28 pieces decrease; exact one-sided sign-test p = 3.73 x 10^-9.")


def s10_revision(prs: Presentation) -> None:
    slide = _content_slide(prs, 10, "Result", "Why the measurement revision matters")
    _stack(
        slide,
        [
            (
                "Initial null",
                "Constant delay plus lag-tolerant envelope coupling showed little effect.",
            ),
            (
                "Control diagnosis",
                "The metric could absorb the lag, so it did not measure simultaneous note attacks.",
            ),
            (
                "Operational revision",
                "Zero-lag onset synchrony directly measures whether singers land attacks together.",
            ),
            (
                "Dissociation",
                "Timing falls 56% to 75%; pure envelope coupling falls roughly 6% to 14% across corpora.",
            ),
        ],
        y=1.45,
        height=1.0,
    )
    _takeaway(
        slide, "The hypothesis stayed fixed; the dated metric revision is part of the audit trail."
    )


def s11_h2(prs: Presentation) -> None:
    slide = _content_slide(prs, 11, "Result", "H2: limited leadership evidence")
    stats = [
        ("0.162", "Observed mean"),
        ("0.155", "Matched null"),
        ("1/5 - 1/3", "Dagstuhl - ESMUC"),
        ("0/20", "ChoralSynth"),
    ]
    x = MARGIN
    for number, label in stats:
        _stat_card(slide, x, Inches(1.35), Inches(2.9), Inches(1.25), number, label, number_size=25)
        x += Inches(3.05)
    _picture_fit(
        slide,
        FIG_DIR / "wp3_flagship_LI_QuartetA_Take02_standard.png",
        MARGIN,
        Inches(2.85),
        SLIDE_W - 2 * MARGIN,
        Inches(3.45),
    )
    _takeaway(slide, "Two human pieces are significant; corpus-level evidence is not significant.")


def s12_h3(prs: Presentation) -> None:
    slide = _content_slide(prs, 12, "Result", "H3: an informative exploratory null")
    stats = [
        ("17 / 18", "videos analyzable"),
        ("1 / 17", "significant, consistent with chance"),
        ("0.068", "median maximum-lag r"),
    ]
    x = MARGIN
    for number, label in stats:
        _stat_card(
            slide, x, Inches(1.35), Inches(3.95), Inches(1.25), number, label, number_size=25
        )
        x += Inches(4.1)
    _picture_fit(
        slide, FIG_DIR / "h3_visual_onset.png", MARGIN, Inches(2.85), Inches(7.35), Inches(3.5)
    )
    _card(
        slide,
        Inches(8.15),
        Inches(2.95),
        Inches(4.55),
        Inches(3.2),
        "What was tested",
        [
            "Maximum signed correlation within 2 s over the first 60-72 s pose window, with 1,000 circular-shift nulls. Full H3 still needs synchronized per-singer audio and video."
        ],
        size=13.2,
    )
    _takeaway(
        slide, "The mixed-audio substitute is null; the paired-data requirement is demonstrated."
    )


def s13_demo(prs: Presentation) -> None:
    slide = _content_slide(prs, 13, "Demo", "Live dashboard demonstration - 60 seconds")
    _stack(
        slide,
        [
            ("1. Audio and network", "Dagstuhl piece: E(t) timeline and influence graph."),
            ("2. Video and pose", "Tier-1 video: playback with pose overlay."),
            (
                "3. Signal availability",
                "Metadata shows which channels each piece actually contains.",
            ),
        ],
        y=1.65,
        height=1.1,
        gap=0.22,
        size=14,
    )
    _takeaway(slide, "The dashboard uses the same envelope-only E(t) definition as the report.")


def s14_limits(prs: Presentation) -> None:
    slide = _content_slide(prs, 14, "Limits", "Weaknesses of the current evidence")
    _stack(
        slide,
        [
            ("Simulated latency", "Models transmission, not live singer adaptation."),
            (
                "Co-varying regimes",
                "Delay, jitter, and dropout change together; individual effects are not isolated.",
            ),
            ("Multimodal gap", "No synchronized per-singer audio and video item exists."),
            ("H2 strength", "Two individual results, but no corpus-level significance."),
            (
                "External validity",
                "E(t) is not yet validated against perceived performance quality.",
            ),
        ],
        y=1.35,
        height=0.85,
        gap=0.12,
        size=12.5,
    )


def s15_results(prs: Presentation) -> None:
    slide = _content_slide(prs, 15, "Results", "Main results and implementation")
    _stack(
        slide,
        [
            (
                "H1 supported for timing",
                "All 28 pieces lose onset synchrony; corpus mean drop 70.7%.",
            ),
            (
                "H2 limited",
                "Two human pieces show significant centralization; overall evidence remains weak.",
            ),
            (
                "H3 exploratory null",
                "Mixed audio plus one pose stream does not substitute for paired per-singer data.",
            ),
            (
                "System delivered",
                "Reproducible pipelines, tested dashboard, figures, report, and presentation package.",
            ),
        ],
        y=1.5,
        height=1.0,
    )


def s16_extensions(prs: Presentation) -> None:
    slide = _content_slide(prs, 16, "Extensions", "What the next study should add")
    _stack(
        slide,
        [
            (
                "Live latency sessions",
                "Measure behavioral adaptation and latency-driven leadership.",
            ),
            ("Paired audio-video corpus", "Unlock the planned H3 incremental-validity test."),
            ("Factorial network experiment", "Separate delay, jitter, and dropout effects."),
            (
                "External validation",
                "Compare E(t) with expert ratings, singer experience, and score-alignment errors.",
            ),
        ],
        y=1.5,
        height=1.0,
    )


def s17_retrospective(prs: Presentation) -> None:
    slide = _content_slide(prs, 17, "Retrospective", "What worked and what could improve")
    _card(
        slide,
        MARGIN,
        Inches(1.45),
        Inches(5.85),
        Inches(4.55),
        "What worked well",
        [
            "Controls exposed the first metric mismatch.\n\nSeeded scripts, tests, and audits kept claims traceable.\n\nInstructor questions improved definitions and limitation framing."
        ],
        size=15,
    )
    _card(
        slide,
        Inches(6.85),
        Inches(1.45),
        Inches(5.85),
        Inches(4.55),
        "What could improve",
        [
            "Secure paired multimodal data earlier.\n\nResolve work-package dependencies before parallel implementation.\n\nConfirm final format and metric expectations earlier."
        ],
        size=15,
    )
    _text(
        slide,
        MARGIN,
        Inches(6.25),
        SLIDE_W - 2 * MARGIN,
        Inches(0.45),
        "Thank you. Questions are welcome.",
        size=22,
        bold=True,
        color=TEAL,
        align=PP_ALIGN.CENTER,
    )


def s18_dashboard_backup(prs: Presentation) -> None:
    slide = _content_slide(prs, 18, "Backup", "Dashboard on real outputs")
    _picture_fit(
        slide,
        FIG_DIR / "wp4_dashboard_realdata.png",
        MARGIN,
        Inches(1.4),
        SLIDE_W - 2 * MARGIN,
        Inches(5.2),
    )
    _takeaway(slide, "Use only if the live dashboard cannot run.")


def s19_repro_backup(prs: Presentation) -> None:
    slide = _content_slide(prs, 19, "Backup", "Reproducibility protocol")
    _stack(
        slide,
        [
            ("Regenerate", "`make reproduce` rebuilds report-stage results and the final PDF."),
            ("H1", "Exact paired sign test plus seeded bootstrap."),
            ("H2", "1,000 matched random graphs per piece with plus-one p-values."),
            ("H3", "1,000 circular shifts per video with deterministic seeds."),
            ("Environment", "Python 3.11.9 and locked dependencies."),
        ],
        y=1.35,
        height=0.85,
        gap=0.12,
        size=12.5,
    )


def build() -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    for builder in (
        s1_title,
        s2_goals,
        s3_related,
        s4_structure,
        s5_data,
        s6_method,
        s7_process,
        s8_repro,
        s9_h1,
        s10_revision,
        s11_h2,
        s12_h3,
        s13_demo,
        s14_limits,
        s15_results,
        s16_extensions,
        s17_retrospective,
        s18_dashboard_backup,
        s19_repro_backup,
    ):
        builder(prs)
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    print(f"Wrote {OUT_PATH} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
